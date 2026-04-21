"""Document extraction and chunking for RAG ingestion.

Reads PDF/DOCX/PPTX/TXT content (from bytes or GCS), extracts text, and returns
fixed-size character chunks. Used by IngestionPipeline for source_type=documents.

Blobs whose name starts with ``_pre_xlsx_`` are treated as preprocessed spreadsheet
text: table-aware chunking (``PRE_XLSX_MAX_CHARS_PER_CHUNK``, typically 3000 chars; no overlap;
strategy-4 row caps).
"""

from __future__ import annotations

from io import BytesIO

from src.services.rag_engine.ingestion.pre_xlsx_documents import (
    PRE_XLSX_MAX_CHARS_PER_CHUNK,
    PreXlsxDocumentsChunker,
)
from src.utils.logger import get_logger


logger = get_logger(__name__)


class DocumentsChunker:
    """Extract text from document bytes and chunk by character size with overlap."""

    def __init__(
        self,
        chunk_size: int = 1500,
        overlap: int = 300,
    ) -> None:
        self.chunk_size = chunk_size
        self.overlap = overlap
        self._pre_xlsx_chunker = PreXlsxDocumentsChunker(
            max_chars_per_chunk=PRE_XLSX_MAX_CHARS_PER_CHUNK,
        )

    def extract_text(self, content: bytes, blob_name: str) -> str:
        """Extract plain text from PDF, DOCX, PPTX, or TXT bytes.

        Args:
            content: Raw file bytes.
            blob_name: Object name (used for extension; e.g. "report.pdf").

        Returns:
            Extracted text as a single string.

        Raises:
            ValueError: If file type is not PDF, DOCX, PPTX, or TXT.
        """
        ext = (blob_name or "").lower()
        if ext.endswith((".txt", ".text")):
            return self._extract_txt(content)
        if ext.endswith(".pdf"):
            return self._extract_pdf(content)
        if ext.endswith(".docx"):
            return self._extract_docx(content)
        if ext.endswith(".pptx"):
            return self._extract_pptx(content)
        raise ValueError("Unsupported file type. Only PDF, DOCX, PPTX, TXT supported.")

    def _extract_pdf(self, content: bytes) -> str:
        import logging

        from pypdf import PdfReader

        # Suppress pypdf FloatObject warnings (invalid numbers in PDF); extraction continues.
        logging.getLogger("pypdf").setLevel(logging.ERROR)

        reader = PdfReader(BytesIO(content))
        parts = []
        for page in reader.pages:
            try:
                t = page.extract_text()
                if t:
                    parts.append(t)
            except Exception:
                logger.warning("Failed to extract text from PDF page", exc_info=True)
                continue
        return " ".join(parts)

    def _extract_docx(self, content: bytes) -> str:
        from docx import Document

        doc = Document(BytesIO(content))
        return "\n\n".join(
            p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()
        )

    def _extract_txt(self, content: bytes) -> str:
        """Decode plain text bytes as UTF-8."""
        return content.decode("utf-8", errors="replace")

    def _extract_pptx(self, content: bytes) -> str:
        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            parts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape.text.strip()
            ]
            if parts:
                slides_text.append(f"Slide {i + 1}:\n" + "\n".join(parts))
        return "\n\n".join(slides_text)

    def chunk_text(self, text: str) -> list[str]:
        """Split text into overlapping character chunks.

        Args:
            text: Full extracted text.

        Returns:
            List of chunk strings.
        """
        if not (text or text.strip()):
            return []
        step = self.chunk_size - self.overlap
        return [text[i : i + self.chunk_size] for i in range(0, len(text), step)]

    def extract_and_chunk(
        self, content: bytes, blob_name: str, opportunity_id: str = ""
    ) -> list[str]:
        """Extract text from document bytes and return chunk list.

        Args:
            content: Raw file bytes (PDF, DOCX, PPTX, or TXT).
            blob_name: Object name for extension detection (e.g. "report.pdf").
            opportunity_id: Opportunity ID for log correlation.

        Returns:
            List of chunk strings for embedding; empty if no text extracted.
        """
        logger.bind(opportunity_id=opportunity_id).info(
            "Documents chunker: extracting and chunking"
        )
        if PreXlsxDocumentsChunker.is_pre_xlsx_blob(blob_name):
            return self._pre_xlsx_chunker.chunk(content, blob_name, opportunity_id)
        raw_text = self.extract_text(content, blob_name)
        chunks = self.chunk_text(raw_text)
        logger.bind(opportunity_id=opportunity_id).info(
            "Documents chunker: extracted and chunked — chunk_count=%s",
            len(chunks),
        )
        return chunks
