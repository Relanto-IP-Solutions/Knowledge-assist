"""Native text extraction for DOCX, MD, PPTX, and XLSX — no LLM calls."""

from __future__ import annotations

from io import BytesIO

from src.utils.logger import get_logger


logger = get_logger(__name__)


class NativeExtractor:
    """Extract text from DOCX, Markdown, PPTX, and XLSX using native libraries."""

    def extract(self, content: bytes, object_name: str) -> str:
        """Extract text from document bytes based on file extension.

        Args:
            content: Raw file bytes.
            object_name: Object name for extension detection (e.g. report.docx).

        Returns:
            Extracted text as a string.

        Raises:
            ValueError: If file type is not DOCX, MD, PPTX, or XLSX.
        """
        ext = (object_name or "").lower()
        if ext.endswith(".docx"):
            return self._extract_docx(content)
        if ext.endswith((".md", ".markdown")):
            return self._extract_md(content)
        if ext.endswith(".pptx"):
            return self._extract_pptx(content)
        if ext.endswith(".xlsx"):
            return self._extract_xlsx(content)
        raise ValueError(
            f"Unsupported file type for native extraction. "
            f"Expected .docx, .md, .pptx, or .xlsx; got: {object_name}"
        )

    def _extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX bytes."""
        from docx import Document

        doc = Document(BytesIO(content))
        return "\n\n".join(
            p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()
        )

    def _extract_md(self, content: bytes) -> str:
        """Extract text from Markdown — decode as UTF-8."""
        return content.decode("utf-8")

    def _extract_pptx(self, content: bytes) -> str:
        """Extract text from PPTX bytes."""
        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            parts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape.text.strip()
            ]
            if parts:
                slides_text.append(f"Slide {i + 1}:\n" + "\n".join(parts))
        return "\n\n".join(slides_text)

    def _extract_xlsx(self, content: bytes) -> str:
        """Extract text from XLSX bytes: sheet names and cell values as plain text."""
        from openpyxl import load_workbook

        wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            sheet_parts = [f"Sheet: {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                row_strs = [
                    str(cell).strip()
                    for cell in row
                    if cell is not None and str(cell).strip()
                ]
                if row_strs:
                    sheet_parts.append(" | ".join(row_strs))
            if len(sheet_parts) > 1:
                parts.append("\n".join(sheet_parts))
        wb.close()
        return "\n\n".join(parts) if parts else ""
